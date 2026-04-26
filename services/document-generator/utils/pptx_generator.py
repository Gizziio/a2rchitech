from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from typing import List, Optional
import tempfile
import os


def create_photo_card_deck(
    title: str,
    key_points: List[str],
    slide_count: int,
    citations: List[dict],
    images: Optional[List[str]] = None
) -> str:
    """
    Create a photo-card style PPTX deck with the given parameters.
    
    Args:
        title: Deck title
        key_points: List of key points (one per slide)
        slide_count: Number of slides to create
        citations: List of citations to include
        images: Optional list of image URLs to use as backgrounds
        
    Returns:
        Path to the generated PPTX file
    """
    prs = Presentation()
    
    # Set slide dimensions (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    subtitle_placeholder.text = "Photo-Card Style Deck"
    
    # Content slides
    content_slide_layout = prs.slide_layouts[6]  # Blank layout for photo cards
    
    for i in range(min(slide_count, len(key_points))):
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Split key point into headline and subline
        key_point = key_points[i]
        parts = key_point.split(" - ", 1)
        if len(parts) == 1:
            parts = key_point.split(": ", 1)
        if len(parts) == 1:
            parts = [key_point[:50] + "..." if len(key_point) > 50 else key_point, ""]
        
        headline = parts[0][:50]  # Limit headline to 50 chars
        subline = parts[1] if len(parts) > 1 else ""
        
        # Add headline (top of slide)
        headline_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
        headline_frame = headline_box.text_frame
        headline_frame.text = headline
        headline_frame.paragraphs[0].font.size = Pt(32)
        headline_frame.paragraphs[0].font.bold = True
        headline_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        headline_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add subline (center of slide)
        subline_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.33), Inches(1.5))
        subline_frame = subline_box.text_frame
        subline_frame.text = subline
        subline_frame.paragraphs[0].font.size = Pt(18)
        subline_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        subline_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add a simple border/frame
        left = top = Inches(0.25)
        width = Inches(12.83)
        height = Inches(7)
        shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 is for rectangle
        shape.line.color.rgb = RGBColor(200, 200, 200)
        shape.fill.background()
    
    # Add citations slide if there are citations
    if citations:
        citation_slide = prs.slides.add_slide(content_slide_layout)
        
        # Add title for citations slide
        title_box = citation_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Sources"
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add citations
        y_pos = Inches(1.5)
        for i, citation in enumerate(citations):
            if i >= 5:  # Limit to 5 citations to avoid overcrowding
                break
            
            citation_text = f"• {citation['title']}\n  {citation['url']}"
            citation_box = citation_slide.shapes.add_textbox(Inches(1), y_pos, Inches(11.33), Inches(1))
            citation_frame = citation_box.text_frame
            citation_frame.text = citation_text
            citation_frame.paragraphs[0].font.size = Pt(14)
            citation_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)  # Dark blue for links
            citation_frame.word_wrap = True
            
            y_pos += Inches(1.2)
    
    # Save to a temporary file
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(temp_file.name)
    temp_file.close()
    
    return temp_file.name